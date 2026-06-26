# -*- coding: utf-8 -*-
"""TripCraft 최종 발표 덱(.pptx) 빌드 — 발표_클로드디자인.md 사양 기반.
   다이어그램은 assets/*.png(render.py 산출물)을 임베드."""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = pathlib.Path(__file__).parent
ASSETS = HERE / "assets"
FONT = "Apple SD Gothic Neo"

PURPLE = RGBColor(0x6D, 0x28, 0xD9)
DPURPLE = RGBColor(0x4C, 0x1D, 0x95)
ACCENT = RGBColor(0xA7, 0x8B, 0xFA)
LILAC = RGBColor(0xED, 0xE9, 0xFE)
TEXT = RGBColor(0x1F, 0x29, 0x37)
SUB = RGBColor(0x6B, 0x72, 0x80)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LBG = RGBColor(0xFA, 0xF9, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

_page = 0


def set_font(run, name=FONT, size=None, bold=None, color=None):
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_runs(p, text, size, color=TEXT, bold=False):
    """**...** 구간을 bold 처리해 runs 추가."""
    parts = text.split("**")
    for i, part in enumerate(parts):
        if part == "":
            continue
        r = p.add_run()
        r.text = part
        set_font(r, size=size, bold=(bold or i % 2 == 1), color=(PURPLE if i % 2 == 1 else color))


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tf


def rect(slide, l, t, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def new_slide(bg=WHITE):
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def footer(slide):
    global _page
    _page += 1
    tf = textbox(slide, Inches(0.5), Inches(7.05), Inches(3), Inches(0.35))
    r = tf.paragraphs[0].add_run(); r.text = "TripCraft"
    set_font(r, size=10, bold=True, color=ACCENT)
    tf2 = textbox(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.35))
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r2 = tf2.paragraphs[0].add_run(); r2.text = str(_page)
    set_font(r2, size=10, color=SUB)


def title(slide, text):
    tf = textbox(slide, Inches(0.7), Inches(0.55), Inches(12), Inches(0.9))
    r = tf.paragraphs[0].add_run(); r.text = text
    set_font(r, size=30, bold=True, color=PURPLE)
    rect(slide, Inches(0.72), Inches(1.45), Inches(1.4), Pt(4), fill=PURPLE)


def body_top():
    return Inches(1.75)


# ---------- 슬라이드 타입 빌더 ----------

def s_cover(d):
    slide = new_slide(LBG)
    rect(slide, 0, 0, SW, Inches(0.25), fill=PURPLE)
    rect(slide, 0, Inches(7.25), SW, Inches(0.25), fill=PURPLE)
    tf = textbox(slide, Inches(1), Inches(2.4), Inches(11.3), Inches(1.3), MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = "TripCraft"
    set_font(r, size=66, bold=True, color=PURPLE)
    tf2 = textbox(slide, Inches(1), Inches(3.7), Inches(11.3), Inches(0.8))
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf2.paragraphs[0].add_run(); r.text = d["sub"]
    set_font(r, size=20, color=TEXT)
    tf3 = textbox(slide, Inches(1), Inches(4.7), Inches(11.3), Inches(0.9))
    for line, sz, col in [(d["team"], 18, DPURPLE), (d["date"], 15, SUB)]:
        p = tf3.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line; set_font(r, size=sz, bold=(col == DPURPLE), color=col)


def s_divider(d):
    slide = new_slide(PURPLE)
    tf = textbox(slide, Inches(1), Inches(2.6), Inches(11.3), Inches(2), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = d["num"]
    set_font(r, size=80, bold=True, color=ACCENT)
    p2 = tf.add_paragraph()
    r = p2.add_run(); r.text = d["title"]
    set_font(r, size=40, bold=True, color=WHITE)


def s_toc(d):
    slide = new_slide()
    title(slide, "목차")
    tf = textbox(slide, Inches(0.9), body_top(), Inches(11.5), Inches(4.8))
    for i, item in enumerate(d["items"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        add_runs(p, item, 20)
    footer(slide)


def s_bullets(d):
    slide = new_slide()
    title(slide, d["title"])
    top = body_top()
    if d.get("lead"):
        tf0 = textbox(slide, Inches(0.7), top, Inches(12), Inches(0.7))
        add_runs(tf0.paragraphs[0], d["lead"], 17, color=DPURPLE)
        top = Inches(2.55)
    tf = textbox(slide, Inches(0.85), top, Inches(11.7), Inches(4.4))
    for i, item in enumerate(d["items"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        b = p.add_run(); b.text = "•  "; set_font(b, size=19, bold=True, color=PURPLE)
        add_runs(p, item, 18)
    if d.get("note"):
        nb = rect(slide, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.55), fill=LILAC)
        ntf = nb.text_frame; ntf.word_wrap = True
        ntf.margin_left = Pt(10); ntf.margin_top = Pt(4)
        add_runs(ntf.paragraphs[0], d["note"], 13, color=DPURPLE)
    footer(slide)


def s_cards(d):
    slide = new_slide()
    title(slide, d["title"])
    cards = d["cards"]
    n = len(cards)
    cols = 2 if n <= 4 else 3
    import math
    rows = math.ceil(n / cols)
    gap = Inches(0.3)
    cw = (Inches(12) - gap * (cols - 1)) / cols
    ch = Inches(1.5) if rows <= 2 else Inches(1.15)
    x0, y0 = Inches(0.7), body_top()
    for i, (head, bodytxt) in enumerate(cards):
        r, c = divmod(i, cols)
        l = x0 + (cw + gap) * c
        t = y0 + (ch + gap) * r
        card = rect(slide, l, t, cw, ch, fill=LBG, line=ACCENT, line_w=Pt(1.25), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = card.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(12); tf.margin_top = Pt(10)
        p = tf.paragraphs[0]; rr = p.add_run(); rr.text = head
        set_font(rr, size=16, bold=True, color=PURPLE)
        p2 = tf.add_paragraph(); p2.space_before = Pt(4)
        add_runs(p2, bodytxt, 13)
    footer(slide)


def s_table(d):
    slide = new_slide()
    title(slide, d["title"])
    top = body_top()
    if d.get("lead"):
        tf0 = textbox(slide, Inches(0.7), top, Inches(12), Inches(0.7))
        add_runs(tf0.paragraphs[0], d["lead"], 16, color=DPURPLE)
        top = top + Inches(0.75)
    headers = d["headers"]; rows = d["rows"]
    nrows = len(rows) + 1; ncols = len(headers)
    th = min(Inches(0.5) * nrows, Inches(4.6))
    gtbl = slide.shapes.add_table(nrows, ncols, Inches(0.7), top, Inches(11.9), th).table
    widths = d.get("widths")
    if widths:
        total = sum(widths)
        for c, wv in enumerate(widths):
            gtbl.columns[c].width = Emu(int(Inches(11.9) * wv / total))
    for c, htext in enumerate(headers):
        cell = gtbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE
        cell.margin_top = cell.margin_bottom = Pt(3)
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = htext; set_font(r, size=14, bold=True, color=WHITE)
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gtbl.cell(ri, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LBG
            cell.margin_left = Pt(8); cell.margin_top = cell.margin_bottom = Pt(2)
            p = cell.text_frame.paragraphs[0]
            if c == 0:
                add_runs(p, val, 12.5, color=DPURPLE, bold=True)
            else:
                add_runs(p, val, 12.5)
    if d.get("note"):
        ntf = textbox(slide, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.5))
        add_runs(ntf.paragraphs[0], d["note"], 12.5, color=SUB)
    footer(slide)


def _fit(img_path, box_w, box_h):
    w, h = Image.open(img_path).size
    ratio = min(box_w / w, box_h / h)
    return int(w * ratio), int(h * ratio)


def s_diagram(d):
    slide = new_slide()
    title(slide, d["title"])
    img = ASSETS / (d["img"] + ".png")
    note_h = Inches(0.5) if d.get("note") else Inches(0)
    box_l, box_t = Inches(0.7), body_top()
    box_w, box_h = Inches(11.9), Inches(4.9) - note_h
    panel = rect(slide, box_l, box_t, box_w, box_h, fill=LBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if img.exists():
        fw, fh = _fit(str(img), int(box_w) - int(Inches(0.4)), int(box_h) - int(Inches(0.4)))
        pic_l = box_l + (box_w - fw) // 2
        pic_t = box_t + (box_h - fh) // 2
        slide.shapes.add_picture(str(img), pic_l, pic_t, width=fw, height=fh)
    else:
        tf = panel.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = f"[다이어그램: {d['img']}]"; set_font(r, size=16, color=SUB)
    if d.get("note"):
        ntf = textbox(slide, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.5))
        add_runs(ntf.paragraphs[0], d["note"], 13, color=DPURPLE)
    footer(slide)


def s_ascii(d):
    slide = new_slide()
    title(slide, d["title"])
    if d.get("items"):
        tf = textbox(slide, Inches(0.85), body_top(), Inches(11.7), Inches(1.3))
        for i, item in enumerate(d["items"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            b = p.add_run(); b.text = "•  "; set_font(b, size=16, bold=True, color=PURPLE)
            add_runs(p, item, 15)
        atop = Inches(3.2)
    else:
        atop = body_top()
    box = rect(slide, Inches(0.7), atop, Inches(11.9), Inches(3.5), fill=RGBColor(0x1F, 0x1B, 0x2E), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = box.text_frame; tf.word_wrap = False
    tf.margin_left = Pt(14); tf.margin_top = Pt(10)
    for i, line in enumerate(d["ascii"].split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line if line else " "
        set_font(r, name="Menlo", size=11.5, color=RGBColor(0xE9, 0xE2, 0xFF))
    footer(slide)


def s_twocol(d):
    slide = new_slide()
    title(slide, d["title"])
    top = body_top()
    for idx, col in enumerate((d["left"], d["right"])):
        l = Inches(0.7) + idx * Inches(6.15)
        panel = rect(slide, l, top, Inches(5.85), Inches(4.6), fill=LBG, line=ACCENT, line_w=Pt(1.25), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = panel.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(14); tf.margin_top = Pt(12)
        p = tf.paragraphs[0]; r = p.add_run(); r.text = col["head"]
        set_font(r, size=18, bold=True, color=PURPLE)
        for item in col["items"]:
            pp = tf.add_paragraph(); pp.space_before = Pt(7)
            b = pp.add_run(); b.text = "• "; set_font(b, size=14, bold=True, color=ACCENT)
            add_runs(pp, item, 14)
    footer(slide)


def s_placeholder(d):
    slide = new_slide()
    title(slide, d["title"])
    boxes = d["boxes"]
    n = len(boxes)
    top = body_top()
    if n == 1:
        ph = rect(slide, Inches(2.2), top + Inches(0.2), Inches(8.9), Inches(4.3), fill=RGBColor(0xF3, 0xF1, 0xFB), line=ACCENT, line_w=Pt(1.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = ph.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = boxes[0]; set_font(r, size=20, bold=True, color=PURPLE)
    else:
        cols = min(n, 3 if n != 4 else 2)
        import math
        rows = math.ceil(n / cols)
        gap = Inches(0.3)
        cw = (Inches(12) - gap * (cols - 1)) / cols
        chh = (Inches(4.6) - gap * (rows - 1)) / rows
        for i, label in enumerate(boxes):
            r, c = divmod(i, cols)
            l = Inches(0.7) + (cw + gap) * c
            t = top + (chh + gap) * r
            ph = rect(slide, l, t, cw, chh, fill=RGBColor(0xF3, 0xF1, 0xFB), line=ACCENT, line_w=Pt(1.25), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            tf = ph.text_frame; tf.word_wrap = True
            pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
            rr = pp.add_run(); rr.text = label; set_font(rr, size=14, bold=True, color=DPURPLE)
    if d.get("note"):
        ntf = textbox(slide, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.6))
        add_runs(ntf.paragraphs[0], d["note"], 12.5, color=SUB)
    footer(slide)


def s_closing(d):
    slide = new_slide(PURPLE)
    tf = textbox(slide, Inches(1), Inches(2.7), Inches(11.3), Inches(2), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "감사합니다"; set_font(r, size=54, bold=True, color=WHITE)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run(); r.text = d["sub"]; set_font(r, size=20, color=LILAC)


BUILDERS = {
    "cover": s_cover, "divider": s_divider, "toc": s_toc, "bullets": s_bullets,
    "cards": s_cards, "table": s_table, "diagram": s_diagram, "ascii": s_ascii,
    "twocol": s_twocol, "placeholder": s_placeholder, "closing": s_closing,
}

from slides_data import SLIDES  # noqa: E402

# 목차번호 자동 부여: 디바이더(또는 sec 키)가 섹션을 열고, 이후 콘텐츠 슬라이드 제목에 "N-n" 접두.
NUMBERED = {"bullets", "table", "cards", "diagram", "ascii", "twocol", "placeholder"}
cur_sec = None
sub = 0
for d in SLIDES:
    if d["type"] == "divider":
        num = d["num"]
        cur_sec = str(int(num)) if num.isdigit() else num
        sub = 0
    elif d.get("sec"):
        cur_sec = d["sec"]
        sub = 0
    if d["type"] in NUMBERED and cur_sec:
        sub += 1
        d["title"] = f"{cur_sec}-{sub}  " + d["title"]
    BUILDERS[d["type"]](d)

out = HERE / "발표.pptx"
prs.save(str(out))
print(f"saved {out} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
